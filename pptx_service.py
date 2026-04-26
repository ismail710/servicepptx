"""
pptx_service.py
---------------
Standalone HTTP service that fills template1.pptx with project data
and returns the result as a base64-encoded PPTX.

Deploy anywhere:
  - Local:      python pptx_service.py
  - Render.com: uses render.yaml in this repo (free tier)
  - Railway:    uses Procfile in this repo (free tier)

POST /generate
  Body (JSON):
    { "project": { ...fields... } }
  Response (JSON):
    { "pptxBase64": "...", "fileName": "PROJECT-001_Title_20260426.pptx" }

GET /health  → 200 { "status": "ok" }
"""

import base64
import io
import json
import os
import re
import sys
import traceback
from datetime import datetime
from http.server import BaseHTTPRequestHandler, HTTPServer

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Pt

# ---------------------------------------------------------------------------
# Path to the PPTX template (same folder as this script by default)
# ---------------------------------------------------------------------------
_SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PATH = os.environ.get("TEMPLATE_PATH", os.path.join(_SCRIPT_DIR, "template1.pptx"))
FIT_FONT_FAMILY = os.environ.get("PPTX_FONT_FAMILY", "Liberation Sans")

SLIDE_ONE_SHAPES = {
    "strategic_axis": 5,
    "closing_date": 6,
    "starting_date": 8,
    "objective": 10,
    "budget_amount": 14,
    "risk_alert": 16,
    "finality": 21,
    "partners": 23,
    "owner_box": 27,
    "co_owner_box": 26,
    "founder_box": 28,
    "capex_box": 54,
    "opex_box": 53,
    "key_milestones": 4,
    "exploit_label": 33,
    "explore_label": 39,
    "capex_commitments": 18,
    "capex_expenses": 30,
    "deliverables": 43,
    "tasks_done": 87,
    "tasks_in_progress": 89,
    "tasks_next_steps": 88,
    "opex_commitments": 48,
    "opex_expenses": 50,
    "header": 75,
}
DEFAULT_MARKER_SHAPES = (52, 55, 56)
MARK_COLOR = RGBColor(147, 112, 10)


def _project_value(project_data: dict, *keys: str) -> str:
    for key in keys:
        value = project_data.get(key)
        if value is not None and str(value).strip() != "":
            return str(value).strip()
    return ""


def _split_items(value) -> list[str]:
    if value is None:
        return []
    if isinstance(value, list):
        raw_items = value
    else:
        text = str(value).replace("\r", "\n")
        if "|" in text:
            raw_items = text.split("|")
        else:
            raw_items = text.split("\n")
    return [item.strip(" \t\n•") for item in raw_items if item and item.strip(" \t\n•")]


def _is_truthy(value) -> bool:
    return str(value).strip().lower() in {"1", "true", "yes", "y", "oui", "x", "checked"}

# ---------------------------------------------------------------------------
# Placeholder → project field mapping
# ---------------------------------------------------------------------------
PLACEHOLDER_MAP = {
    "{{Project_Number}}":     lambda p: p.get("projectNumber", ""),
    "{{Project_Title}}":      lambda p: p.get("projectTitle", ""),
    "{{Project_Lead}}":       lambda p: p.get("projectLead", ""),
    "{{PI_UM6P}}":            lambda p: p.get("piUm6p", ""),
    "{{Starting_Date}}":      lambda p: p.get("startDate", ""),
    "{{Closing_Date}}":       lambda p: p.get("endDate", ""),
    "{{Budget_Amount}}":      lambda p: _fmt_number(p.get("budget", "")),
    "{{Deliverables_List}}":  lambda p: _bullets(_project_value(p, "deliverables", "deliverablesCsv")),
    "{{Risk_Alert}}":         lambda p: _bullets(_project_value(p, "risks_and_alerts", "risksAndAlertsCsv")),
    "{{Tasks_Done}}":         lambda p: _bullets(_project_value(p, "done", "doneCsv")),
    "{{Tasks_InProgress}}":   lambda p: _bullets(_project_value(p, "in_progress", "inProgressCsv")),
    "{{Tasks_NextSteps}}":    lambda p: _bullets(_project_value(p, "planned", "plannedCsv")),
    "{{CAPEX_Commitments}}":  lambda p: _fmt_number(p.get("capex_commitments", p.get("capexCommitments", ""))),
    "{{CAPEX_Expenses}}":     lambda p: _fmt_number(p.get("capex_expenses", p.get("capexExpenses", ""))),
    "{{OPEX_Commitments}}":   lambda p: _fmt_number(p.get("opex_commitments", p.get("opexCommitments", ""))),
    "{{OPEX_Expenses}}":      lambda p: _fmt_number(p.get("opex_expenses", p.get("opexExpenses", ""))),
}


def _fmt_number(val: str) -> str:
    """Format numeric strings with thousands separator (e.g. 250000 → 250 000)."""
    v = str(val).strip()
    try:
        n = int(float(v))
        return f"{n:,}".replace(",", " ")
    except (ValueError, TypeError):
        return v


def _bullets(csv_val: str) -> str:
    """Convert pipe-separated CSV to bullet lines: 'A | B' → '• A\n• B'."""
    items = _split_items(csv_val)
    if not items:
        return ""
    return "\n".join(f"• {item}" for item in items)


def _fit_text_frame(text_frame, max_size: int, vertical_anchor=MSO_ANCHOR.TOP):
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.vertical_anchor = vertical_anchor
    try:
        text_frame.fit_text(font_family=FIT_FONT_FAMILY, max_size=max_size)
    except Exception:
        # Render/other hosts may not have the chosen font installed. Keep
        # the auto-fit setting so PowerPoint still shrinks text on open.
        pass


def _set_shape_text(shape, text: str, *, max_size: int, bullet: bool = False, vertical_anchor=MSO_ANCHOR.TOP):
    text_frame = shape.text_frame
    text_frame.clear()

    items = _split_items(text) if bullet else [str(text).strip()]
    if not items:
        items = [""]

    for index, item in enumerate(items):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(0)
        run = paragraph.add_run()
        run.text = f"• {item}" if bullet and item else item
        run.font.size = Pt(max_size)

    _fit_text_frame(text_frame, max_size=max_size, vertical_anchor=vertical_anchor)


def _fit_existing_shape(shape, *, max_size: int, vertical_anchor=MSO_ANCHOR.TOP):
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(max_size)
    _fit_text_frame(text_frame, max_size=max_size, vertical_anchor=vertical_anchor)


def _set_checkbox_mark(shape, selected: bool):
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = "X" if selected else ""
    run.font.bold = True
    run.font.size = Pt(14)
    run.font.color.rgb = MARK_COLOR

    _fit_text_frame(text_frame, max_size=14, vertical_anchor=MSO_ANCHOR.MIDDLE)


def _remove_shape(shape):
    element = shape.element
    element.getparent().remove(element)


def _add_positioning_mark(slide, label_shape):
    mark_width = 135157
    mark_height = 192115
    mark_left = label_shape.left + (label_shape.width - mark_width) // 2
    mark_top = label_shape.top + label_shape.height + 23208
    textbox = slide.shapes.add_textbox(mark_left, mark_top, mark_width, mark_height)
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = "X"
    run.font.bold = True
    run.font.size = Pt(16)
    run.font.color.rgb = MARK_COLOR
    _fit_text_frame(text_frame, max_size=16, vertical_anchor=MSO_ANCHOR.MIDDLE)


def _populate_slide_one(slide, project_data: dict):
    shapes = {
        name: slide.shapes[index - 1]
        for name, index in SLIDE_ONE_SHAPES.items()
        if index <= len(slide.shapes)
    }
    default_marker_shapes = [
        slide.shapes[index - 1]
        for index in DEFAULT_MARKER_SHAPES
        if index <= len(slide.shapes)
    ]

    _set_shape_text(
        shapes["objective"],
        _project_value(project_data, "objective", "projectObjective", "projectDescription"),
        max_size=17,
        vertical_anchor=MSO_ANCHOR.TOP,
    )
    _set_shape_text(
        shapes["partners"],
        _project_value(project_data, "partners", "partnersCsv"),
        max_size=15,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    _set_shape_text(
        shapes["strategic_axis"],
        _project_value(project_data, "strategicAxis"),
        max_size=15,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    _set_shape_text(
        shapes["finality"],
        _project_value(project_data, "finality"),
        max_size=15,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    _set_shape_text(
        shapes["key_milestones"],
        _project_value(project_data, "key_milestones", "keyMilestonesCsv"),
        max_size=16,
        bullet=True,
        vertical_anchor=MSO_ANCHOR.TOP,
    )
    _set_shape_text(
        shapes["deliverables"],
        _project_value(project_data, "deliverables", "deliverablesCsv"),
        max_size=15,
        bullet=True,
        vertical_anchor=MSO_ANCHOR.TOP,
    )
    _set_shape_text(
        shapes["risk_alert"],
        _project_value(project_data, "risks_and_alerts", "risksAndAlertsCsv"),
        max_size=14,
        bullet=True,
        vertical_anchor=MSO_ANCHOR.TOP,
    )
    _set_shape_text(
        shapes["tasks_done"],
        _project_value(project_data, "done", "doneCsv"),
        max_size=15,
        bullet=True,
        vertical_anchor=MSO_ANCHOR.TOP,
    )
    _set_shape_text(
        shapes["tasks_in_progress"],
        _project_value(project_data, "in_progress", "inProgressCsv"),
        max_size=15,
        bullet=True,
        vertical_anchor=MSO_ANCHOR.TOP,
    )
    _set_shape_text(
        shapes["tasks_next_steps"],
        _project_value(project_data, "planned", "plannedCsv"),
        max_size=15,
        bullet=True,
        vertical_anchor=MSO_ANCHOR.TOP,
    )

    for name, size in {
        "header": 19,
        "starting_date": 13,
        "closing_date": 13,
        "budget_amount": 13,
        "capex_commitments": 12,
        "capex_expenses": 12,
        "opex_commitments": 12,
        "opex_expenses": 12,
    }.items():
        _fit_existing_shape(shapes[name], max_size=size, vertical_anchor=MSO_ANCHOR.MIDDLE)

    budget_type = _project_value(project_data, "budgetType", "budget_type").lower()
    _set_checkbox_mark(shapes["capex_box"], budget_type == "capex" or _is_truthy(project_data.get("capexSelected")))
    _set_checkbox_mark(shapes["opex_box"], budget_type == "opex" or _is_truthy(project_data.get("opexSelected")))
    _set_checkbox_mark(shapes["owner_box"], _is_truthy(_project_value(project_data, "owner_selected", "ownerSelected")))
    _set_checkbox_mark(shapes["co_owner_box"], _is_truthy(_project_value(project_data, "co_owner_selected", "coOwnerSelected")))
    _set_checkbox_mark(shapes["founder_box"], _is_truthy(_project_value(project_data, "founder_selected", "founderSelected")))

    for marker_shape in default_marker_shapes:
        _remove_shape(marker_shape)

    project_positioning = _project_value(project_data, "projectPositioning").lower()
    exploit_selected = _is_truthy(_project_value(project_data, "exploit_selected", "exploitSelected")) or project_positioning == "exploit"
    explore_selected = _is_truthy(_project_value(project_data, "explore_selected", "exploreSelected")) or project_positioning == "explore"
    if exploit_selected:
        _add_positioning_mark(slide, shapes["exploit_label"])
    if explore_selected:
        _add_positioning_mark(slide, shapes["explore_label"])


# ---------------------------------------------------------------------------
# Core PPTX filling logic
# ---------------------------------------------------------------------------

def _replace_in_paragraph(para, replacements: dict):
    """
    Replace all {{TOKEN}} occurrences in a paragraph.
    Handles tokens that are split across multiple runs by first merging
    all run text, performing the replacement, then writing back into the
    first run and clearing the rest.
    """
    if not para.runs:
        return

    # Merge all run text
    full_text = "".join(r.text for r in para.runs)

    # Check if any replacement is needed
    if not any(token in full_text for token in replacements):
        return

    # Perform all replacements on the merged text
    new_text = full_text
    for token, value in replacements.items():
        new_text = new_text.replace(token, value)

    # Write back: put everything in the first run, clear subsequent runs
    para.runs[0].text = new_text
    for run in para.runs[1:]:
        run.text = ""


def fill_pptx_base64(project_data: dict) -> dict:
    """
    Load template1.pptx, replace all placeholders, and return
    { "pptxBase64": "...", "fileName": "..." }
    """
    prs = Presentation(TEMPLATE_PATH)

    # Build token → value dict
    replacements = {
        token: fn(project_data)
        for token, fn in PLACEHOLDER_MAP.items()
    }

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                _replace_in_paragraph(para, replacements)

    if prs.slides:
        _populate_slide_one(prs.slides[0], project_data)

    # Serialize to bytes
    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    # Build file name
    title_slug = re.sub(r"[^A-Za-z0-9_]", "_",
                        project_data.get("projectTitle", "Projet"))[:40]
    proj_num   = re.sub(r"[^A-Za-z0-9_-]", "_",
                        project_data.get("projectNumber", "XXX"))
        run_stamp  = datetime.now().strftime("%Y%m%d_%H%M%S_%f")[:-3]
        file_name  = f"{proj_num}_{title_slug}_{run_stamp}.pptx"

    return {
        "pptxBase64": base64.b64encode(pptx_bytes).decode("utf-8"),
        "fileName":   file_name,
    }


# ---------------------------------------------------------------------------
# HTTP request handler
# ---------------------------------------------------------------------------

class PPTXHandler(BaseHTTPRequestHandler):

    def log_message(self, fmt, *args):  # suppress default noisy logging
        print(f"[{self.address_string()}] {fmt % args}", file=sys.stderr)

    def _send_json(self, status: int, data: dict):
        body = json.dumps(data, ensure_ascii=False).encode("utf-8")
        self.send_response(status)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.send_header("Content-Length", str(len(body)))
        self.send_header("Access-Control-Allow-Origin", "*")
        self.end_headers()
        self.wfile.write(body)

    def do_OPTIONS(self):
        self.send_response(204)
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Methods", "POST, GET, OPTIONS")
        self.send_header("Access-Control-Allow-Headers", "Content-Type")
        self.end_headers()

    def do_GET(self):
        if self.path.rstrip("/") in ("/health", "/api/health"):
            self._send_json(200, {"status": "ok", "template": TEMPLATE_PATH})
        else:
            self._send_json(404, {"error": "Not found"})

    def do_POST(self):
        if self.path.rstrip("/") not in ("/generate", "/api/generate"):
            self._send_json(404, {"error": "Not found"})
            return

        try:
            length  = int(self.headers.get("Content-Length", 0))
            raw     = self.rfile.read(length) if length else b"{}"
            payload = json.loads(raw.decode("utf-8"))
        except Exception as exc:
            self._send_json(400, {"error": f"Invalid JSON: {exc}"})
            return

        # Accept both { project: {...} } and flat { projectTitle: ... }
        project_data = payload.get("project", payload)

        try:
            result = fill_pptx_base64(project_data)
            self._send_json(200, result)
        except Exception as exc:
            traceback.print_exc()
            self._send_json(500, {"error": str(exc)})


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def run_server(port: int = 7771):
    if not os.path.exists(TEMPLATE_PATH):
        print(f"ERROR: template not found at {TEMPLATE_PATH}", file=sys.stderr)
        sys.exit(1)
    server = HTTPServer(("0.0.0.0", port), PPTXHandler)
    print(f"PPTX service running on port {port} — POST /generate  GET /health")
    print(f"Template: {TEMPLATE_PATH}")
    server.serve_forever()


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 7771))
    run_server(port)
